from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "quy_trinh_du_an_datn.pptx"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


COLORS = {
    "bg": rgb("#F5F1E8"),
    "ink": rgb("#1F2430"),
    "muted": rgb("#5A6473"),
    "accent": rgb("#BF5B39"),
    "accent2": rgb("#2D6A73"),
    "card": rgb("#FFFDF8"),
    "line": rgb("#D9C7A7"),
}


def add_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(8.5), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "DejaVu Sans"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = COLORS["ink"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(8.8), Inches(0.45))
        p2 = sub.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "DejaVu Sans"
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLORS["muted"]


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16,
             color=COLORS["ink"], bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "DejaVu Sans"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_card(slide, x: float, y: float, w: float, h: float, title: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["card"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1.2)
    if title:
        add_text(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.35, size=13, bold=True, color=COLORS["accent2"])
    return shape


def add_picture(slide, path: Path, x: float, y: float, w: float, h: float):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_label_bar(slide, text: str, x: float, y: float, w: float, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.33)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_text(slide, text, x, y + 0.02, w, 0.25, size=10, color=rgb("#FFFFFF"), bold=True, align=PP_ALIGN.CENTER)


def add_process_box(slide, text: str, x: float, y: float, w: float, h: float, fill):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    add_text(slide, text, x + 0.05, y + 0.10, w - 0.1, h - 0.15, size=12, bold=True,
             color=rgb("#FFFFFF"), align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = COLORS["accent2"]
    line.line.width = Pt(2.0)
    line.line.end_arrowhead = True


def get_eval_summary() -> dict:
    path = ROOT / "openroad_docker_lab" / "alphachip_like_eval_summary.json"
    if path.exists():
        return json.loads(path.read_text())
    return {}


def slide_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Quy Trình Thực Hiện Dự Án RL Macro Placement",
              "Từ benchmark MacroPlacement đến huấn luyện RL, đánh giá và kiểm chứng bằng OpenROAD")
    left = ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/CT_Placement.png"
    mid = ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/Ariane133_ORFS.png"
    right = ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/Human_Gridded_Placement.png"
    add_picture(slide, left, 0.55, 1.45, 4.0, 4.55)
    add_picture(slide, mid, 4.7, 1.45, 4.0, 4.55)
    add_picture(slide, right, 8.85, 1.45, 4.0, 4.55)
    add_label_bar(slide, "Baseline Circuit Training / RL", 0.72, 5.72, 3.66, COLORS["accent"])
    add_label_bar(slide, "Placement OpenROAD / ORFS", 4.87, 5.72, 3.66, COLORS["accent2"])
    add_label_bar(slide, "Baseline thủ công theo lưới", 9.02, 5.72, 3.66, COLORS["accent"])


def slide_problem(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "1. Bài toán và benchmark")
    add_card(slide, 0.55, 1.2, 6.15, 2.15, "Thiết kế và biểu diễn bài toán")
    add_picture(slide, ROOT / "MacroPlacement/Docs/CodeElements/images/macro_example.png", 0.8, 1.55, 5.65, 1.4)
    add_text(slide,
             "Ariane133 là benchmark có 133 macro. RL không đặt từng standard cell mà học cách đặt các hard macro lớn trên canvas.",
             0.85, 2.95, 5.5, 0.55, size=11, color=COLORS["muted"])
    add_card(slide, 7.0, 1.2, 5.85, 5.05, "Đồ thị kết nối cho RL")
    add_picture(slide, ROOT / "MacroPlacement/Docs/CodeElements/images/net_model.png", 7.25, 1.58, 5.35, 3.8)
    add_text(slide,
             "Translator của MacroPlacement chuyển netlist / DEF / LEF / metadata sang netlist.pb.txt và initial.plc, là đầu vào trực tiếp cho environment RL.",
             7.25, 5.5, 5.2, 0.6, size=11, color=COLORS["muted"])


def slide_pipeline(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "2. Luồng EDA trước khi vào RL")
    boxes = [
        ("RTL + SDC", 0.7, COLORS["accent"]),
        ("Synthesis", 2.4, COLORS["accent2"]),
        ("DEF / LEF /\nmetadata", 4.1, COLORS["accent"]),
        ("Chuyển đổi định dạng /\ntranslator", 6.1, COLORS["accent2"]),
        ("netlist.pb.txt\ninitial.plc", 8.45, COLORS["accent"]),
    ]
    prev_end = None
    for text, x, fill in boxes:
        add_process_box(slide, text, x, 1.25, 1.45, 0.8, fill)
        if prev_end is not None:
            add_arrow(slide, prev_end, 1.65, x, 1.65)
        prev_end = x + 1.45
    add_card(slide, 0.65, 2.35, 4.35, 3.55, "Baseline placement trước RL")
    add_picture(slide, ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/Innovus_Flow2_Placement.png", 0.9, 2.72, 1.95, 2.55)
    add_picture(slide, ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/Human_Gridded_Placement.png", 2.95, 2.72, 1.8, 2.55)
    add_text(slide,
             "Benchmark đã có sẵn nhiều baseline placement: CMP, CT, manual, OpenROAD RTL-MP. RL học trên dữ liệu đã được chuyển đổi từ flow EDA này.",
             0.9, 5.35, 3.9, 0.45, size=11, color=COLORS["muted"])
    add_card(slide, 5.3, 2.35, 7.55, 3.55, "Grouping / gridding và tạo dữ liệu RL")
    add_picture(slide, ROOT / "MacroPlacement/Docs/CodeElements/images/Gridding Algorithm.png", 5.55, 2.72, 2.55, 2.6)
    add_picture(slide, ROOT / "MacroPlacement/Docs/OurProgress/images/EvaluationFlows.png", 8.3, 2.72, 4.2, 2.6)
    add_text(slide,
             "Sau gridding, macro centers được map lên canvas. Output quan trọng nhất cho RL là cặp file netlist.pb.txt và initial.plc.",
             5.6, 5.35, 6.7, 0.45, size=11, color=COLORS["muted"])


def slide_rl_loop(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "3. Vòng lặp RL macro placement")
    add_process_box(slide, "PlacementCost\nrestore_placement(initial.plc)", 0.85, 1.55, 2.2, 0.9, COLORS["accent2"])
    add_process_box(slide, "Bộ trích xuất đặc trưng\nobservation_for_node()", 3.45, 1.55, 2.2, 0.9, COLORS["accent"])
    add_process_box(slide, "Mô hình Actor-Critic\nchọn action", 6.05, 1.55, 2.0, 0.9, COLORS["accent2"])
    add_process_box(slide, "place_node()\nđặt macro", 8.4, 1.55, 1.8, 0.9, COLORS["accent"])
    add_process_box(slide, "reward / cost\nwirelength, density", 10.55, 1.55, 2.1, 0.9, COLORS["accent2"])
    for a, b in [(3.05, 3.45), (5.65, 6.05), (8.05, 8.4), (10.2, 10.55)]:
        add_arrow(slide, a, 2.0, b, 2.0)
    add_process_box(slide, "rollout buffer", 4.2, 3.25, 1.9, 0.85, COLORS["accent"])
    add_process_box(slide, "Cập nhật PPO\ncompute returns / advantages\noptimizer.step()", 6.65, 3.0, 2.55, 1.35, COLORS["accent2"])
    add_arrow(slide, 11.6, 2.45, 11.6, 3.65)
    add_arrow(slide, 10.95, 3.65, 9.2, 3.65)
    add_arrow(slide, 6.65, 3.65, 6.1, 3.65)
    add_arrow(slide, 7.95, 4.35, 7.95, 5.15)
    add_text(slide, "episode = đặt xong một lượt macro", 6.55, 5.18, 2.8, 0.25, size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_arrow(slide, 7.95, 5.05, 4.6, 5.05)
    add_arrow(slide, 4.6, 5.05, 4.6, 2.45)
    add_text(slide,
             "Sau vài episode (rollout_episodes), PPO mới cập nhật trọng số. Kết quả train là model .pt, history .csv và summary .json.",
             0.95, 5.65, 11.5, 0.55, size=12, color=COLORS["muted"])


def slide_infra(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "4. Hạ tầng thực nghiệm")
    add_card(slide, 0.75, 1.45, 5.7, 4.35, "Server GPU")
    add_process_box(slide, "Git clone DATN", 1.15, 2.0, 1.7, 0.75, COLORS["accent"])
    add_process_box(slide, "Train PPO trên CUDA", 3.15, 2.0, 2.0, 0.75, COLORS["accent2"])
    add_process_box(slide, "Evaluate\nsinh final.plc", 1.15, 3.1, 1.7, 0.9, COLORS["accent2"])
    add_process_box(slide, "Output:\n.pt / .plc / .json", 3.15, 3.1, 2.0, 0.9, COLORS["accent"])
    add_arrow(slide, 2.85, 2.37, 3.15, 2.37)
    add_arrow(slide, 4.15, 2.75, 4.15, 3.1)
    add_arrow(slide, 2.85, 3.55, 3.15, 3.55)
    add_text(slide,
             "Server được dùng cho tính toán nặng: train và eval policy. Không cần GUI OpenROAD tại đây.",
             1.0, 4.55, 4.9, 0.55, size=11, color=COLORS["muted"])
    add_card(slide, 7.0, 1.45, 5.5, 4.35, "Máy local")
    add_process_box(slide, "Nhận kết quả\nbằng scp", 7.45, 2.0, 1.7, 0.75, COLORS["accent"])
    add_process_box(slide, "OpenROAD GUI", 9.45, 2.0, 1.7, 0.75, COLORS["accent2"])
    add_process_box(slide, "So sánh với\ninitial / legalized", 7.45, 3.1, 1.7, 0.9, COLORS["accent2"])
    add_process_box(slide, "Kiểm chứng\nEDA cuối", 9.45, 3.1, 1.7, 0.9, COLORS["accent"],)
    add_arrow(slide, 9.15, 2.37, 9.45, 2.37)
    add_arrow(slide, 10.3, 2.75, 10.3, 3.1)
    add_arrow(slide, 9.15, 3.55, 9.45, 3.55)
    add_text(slide,
             "Máy local phù hợp cho OpenROAD, GUI và phần so sánh placement với benchmark gốc.",
             7.35, 4.55, 4.8, 0.55, size=11, color=COLORS["muted"])


def slide_results(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "5. Kết quả pipeline hiện tại")
    summary = get_eval_summary()
    add_card(slide, 0.75, 1.35, 5.95, 4.65, "Kiểm chứng pipeline")
    add_picture(slide, ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/CT_Placement.png", 1.0, 1.75, 2.2, 2.45)
    add_picture(slide, ROOT / "MacroPlacement/Flows/NanGate45/ariane133/screenshots/Ariane133_ORFS.png", 3.45, 1.75, 2.95, 2.45)
    add_text(slide,
             "Đã hoàn thành đường đi full pipeline: train trên GPU server -> evaluate -> sinh alphachip_like_final.plc -> kéo về local để phân tích.",
             1.0, 4.4, 5.2, 0.55, size=12, color=COLORS["muted"])
    add_card(slide, 7.0, 1.35, 5.55, 4.65, "Số liệu run gần nhất")
    if summary:
        metrics = [
            ("Initial cost", summary.get("initial_cost")),
            ("Final cost", summary.get("cost")),
            ("Best cost", summary.get("best_cost")),
            ("Wirelength", summary.get("wirelength")),
            ("Steps", summary.get("steps")),
            ("Runtime (s)", summary.get("runtime_sec")),
        ]
        y = 1.8
        for name, value in metrics:
            val = f"{value:.6f}" if isinstance(value, float) else str(value)
            add_text(slide, name, 7.35, y, 2.2, 0.25, size=13, bold=True, color=COLORS["accent2"])
            add_text(slide, val, 9.7, y, 2.35, 0.25, size=13, color=COLORS["ink"])
            y += 0.42
        note = "Run nay chu yeu la smoke test: pipeline dung, nhung final cost chua tot hon baseline initial."
    else:
        note = "Chua tim thay file eval summary tai openroad_docker_lab/alphachip_like_eval_summary.json."
    if summary:
        note = "Run này chủ yếu là smoke test: pipeline đúng, nhưng final cost chưa tốt hơn baseline initial."
    else:
        note = "Chưa tìm thấy file eval summary tại openroad_docker_lab/alphachip_like_eval_summary.json."
    add_text(slide, note, 7.35, 4.55, 4.75, 0.65, size=11, color=COLORS["muted"])


def slide_next(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "6. Bước tiếp theo")
    add_process_box(slide, "Tăng episodes /\nrollout_episodes", 1.0, 1.65, 2.25, 1.0, COLORS["accent"])
    add_process_box(slide, "Chạy nhiều seed\nvà gom metrics", 3.7, 1.65, 2.25, 1.0, COLORS["accent2"])
    add_process_box(slide, "So với initial.plc /\nlegalized.plc", 6.4, 1.65, 2.25, 1.0, COLORS["accent"])
    add_process_box(slide, "Convert về Tcl / DEF\ncho OpenROAD", 9.1, 1.65, 2.25, 1.0, COLORS["accent2"])
    for a, b in [(3.25, 3.7), (5.95, 6.4), (8.65, 9.1)]:
        add_arrow(slide, a, 2.15, b, 2.15)
    add_card(slide, 0.95, 3.05, 11.3, 2.45, "Thông điệp chính")
    add_text(slide,
             "Giá trị của dự án hiện tại là đã thông được toàn bộ quy trình: benchmark MacroPlacement -> format RL -> train PPO -> evaluate -> final.plc -> phân tích bằng OpenROAD.",
             1.25, 3.55, 10.7, 0.6, size=17, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_text(slide,
             "Từ đây, công việc chính là nâng chất lượng học: train lâu hơn, so sánh nhiều baseline hơn, và đưa placement RL quay lại flow EDA để kiểm chứng QoR.",
             1.2, 4.45, 10.8, 0.55, size=13, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_pipeline(prs)
    slide_rl_loop(prs)
    slide_infra(prs)
    slide_results(prs)
    slide_next(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
